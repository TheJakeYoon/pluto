"""File loaders for supported insertion formats.

Supported: pdf, txt, md, jpg, jpeg, png, heic, json, html/htm, doc, docx, xls, xlsx, ppt, pptx.

Strategy:
- Text-native formats (txt/md/json/html/csv) → decoded directly.
- Office formats → python-docx / python-pptx / openpyxl.
- PDF → rendered to page images with PyMuPDF, then OCR'd with pytesseract.
  Native PDF text is kept as a fallback if OCR is unavailable or empty.
- Images → normalized to PNG and OCR'd with pytesseract.
"""

from __future__ import annotations

import base64
import io
import json
import logging
import os
import time
from dataclasses import dataclass, field
from typing import Callable, List, Optional, Tuple

logger = logging.getLogger("uvicorn.error")


def _env_int(name: str, default: int) -> int:
    try:
        return int(os.environ.get(name, str(default)))
    except ValueError:
        return default


def _env_float(name: str, default: float) -> float:
    try:
        return float(os.environ.get(name, str(default)))
    except ValueError:
        return default


PDF_OCR_MAX_PAGES = _env_int("PDF_OCR_MAX_PAGES", 250)  # <= 0 means all pages
PDF_OCR_ZOOM = _env_float("PDF_OCR_ZOOM", 1.8)
PDF_VISION_FALLBACK_MAX_PAGES = _env_int("PDF_VISION_FALLBACK_MAX_PAGES", 20)


@dataclass
class DocumentLoadResult:
    """Unified return type for any loader call."""

    text: str = ""
    page_images: List[bytes] = field(default_factory=list)  # PNG bytes per page/frame
    images: List[bytes] = field(default_factory=list)  # extra embedded images (raw PNG)
    metadata: dict = field(default_factory=dict)
    needs_vision: bool = False  # true when caller should run vision OCR to get text
    warnings: List[str] = field(default_factory=list)


# ---------- individual loaders ----------


def _read_text(path: str) -> str:
    with open(path, "rb") as f:
        raw = f.read()
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("utf-8", errors="replace")


def _load_txt_md(path: str) -> DocumentLoadResult:
    return DocumentLoadResult(text=_read_text(path), metadata={"loader": "text"})


def _load_json(path: str) -> DocumentLoadResult:
    raw = _read_text(path)
    try:
        parsed = json.loads(raw)
        pretty = json.dumps(parsed, indent=2, ensure_ascii=False)
        return DocumentLoadResult(text=pretty, metadata={"loader": "json"})
    except json.JSONDecodeError:
        return DocumentLoadResult(text=raw, metadata={"loader": "json-raw"}, warnings=["Invalid JSON; stored raw."])


def _load_html(path: str) -> DocumentLoadResult:
    raw = _read_text(path)
    try:
        from bs4 import BeautifulSoup  # type: ignore

        soup = BeautifulSoup(raw, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        text = soup.get_text("\n", strip=True)
        title = (soup.title.string or "").strip() if soup.title else ""
        return DocumentLoadResult(text=text, metadata={"loader": "html", "title": title})
    except Exception as e:  # pragma: no cover
        return DocumentLoadResult(text=raw, metadata={"loader": "html-raw"}, warnings=[f"bs4 failed: {e}"])


def extract_pdf_text(path: str) -> Tuple[str, dict]:
    """Return ``(text, meta)`` extracted from a PDF using PyMuPDF."""
    import fitz  # type: ignore

    parts: List[str] = []
    with fitz.open(path) as doc:
        meta = {
            "loader": "pdf",
            "page_count": doc.page_count,
            "pdf_title": (doc.metadata or {}).get("title") or "",
        }
        for i, page in enumerate(doc):
            txt = page.get_text("text") or ""
            if txt.strip():
                parts.append(f"[Page {i + 1}]\n{txt.strip()}")
    return ("\n\n".join(parts).strip(), meta)


def pdf_page_images(path: str, *, max_pages: int = PDF_VISION_FALLBACK_MAX_PAGES, zoom: float = PDF_OCR_ZOOM) -> List[bytes]:
    """Render up to ``max_pages`` PDF pages as PNG bytes (for vision OCR)."""
    import fitz  # type: ignore

    out: List[bytes] = []
    with fitz.open(path) as doc:
        pages = doc.page_count if max_pages <= 0 else min(doc.page_count, max_pages)
        matrix = fitz.Matrix(zoom, zoom)
        for i in range(pages):
            page_start = time.perf_counter()
            pix = doc[i].get_pixmap(matrix=matrix, alpha=False)
            out.append(pix.tobytes("png"))
            logger.info(
                "Insertion PDF render page file=%s page=%d/%d image_bytes=%d elapsed_s=%.2f",
                os.path.basename(path),
                i + 1,
                doc.page_count,
                len(out[-1]),
                time.perf_counter() - page_start,
            )
    return out


def _tesseract_ocr_png(png_bytes: bytes) -> str:
    """OCR PNG bytes with pytesseract.

    ``pytesseract`` is a Python wrapper around the native ``tesseract`` binary,
    so callers catch failures and can surface a useful warning to the UI.
    """
    from PIL import Image  # type: ignore
    import pytesseract  # type: ignore

    with Image.open(io.BytesIO(png_bytes)) as img:
        img = img.convert("RGB")
        return (pytesseract.image_to_string(img) or "").strip()


def _tesseract_ocr_pdf_pages(pages_png: List[bytes]) -> str:
    parts: List[str] = []
    for i, png in enumerate(pages_png, start=1):
        text = _tesseract_ocr_png(png)
        if text.strip():
            parts.append(f"[Page {i}]\n{text.strip()}")
    return "\n\n".join(parts).strip()


def _ensure_tesseract_available() -> str:
    import pytesseract  # type: ignore

    return str(pytesseract.get_tesseract_version())


def _tesseract_ocr_pdf_file(
    path: str,
    *,
    max_pages: int,
    zoom: float,
    progress_callback: Optional[Callable[[dict], None]] = None,
) -> tuple[str, int, int]:
    import fitz  # type: ignore

    parts: List[str] = []
    total_chars = 0
    with fitz.open(path) as doc:
        total_pages = doc.page_count
        pages_to_process = total_pages if max_pages <= 0 else min(total_pages, max_pages)
        matrix = fitz.Matrix(zoom, zoom)
        logger.info(
            "Insertion PDF OCR start file=%s pages_total=%d pages_to_process=%d max_pages=%d zoom=%.2f",
            os.path.basename(path),
            total_pages,
            pages_to_process,
            max_pages,
            zoom,
        )
        if pages_to_process < total_pages:
            logger.warning(
                "Insertion PDF OCR truncated file=%s pages_total=%d pages_processed=%d set PDF_OCR_MAX_PAGES=0 to process all pages",
                os.path.basename(path),
                total_pages,
                pages_to_process,
            )
        for i in range(pages_to_process):
            page_start = time.perf_counter()
            pix = doc[i].get_pixmap(matrix=matrix, alpha=False)
            png = pix.tobytes("png")
            text = _tesseract_ocr_png(png)
            page_text = text.strip()
            page_chars = len(page_text)
            total_chars += page_chars
            logger.info(
                "Insertion PDF OCR page file=%s page=%d/%d chars=%d cumulative_chars=%d image_bytes=%d elapsed_s=%.2f",
                os.path.basename(path),
                i + 1,
                total_pages,
                page_chars,
                total_chars,
                len(png),
                time.perf_counter() - page_start,
            )
            if progress_callback is not None:
                progress_callback(
                    {
                        "event": "pdf_ocr_page",
                        "page": i + 1,
                        "page_count": total_pages,
                        "pages_inserted": i + 1,
                        "chars": page_chars,
                        "cumulative_chars": total_chars,
                    }
                )
            if page_text:
                parts.append(f"[Page {i + 1}]\n{page_text}")
    return "\n\n".join(parts).strip(), pages_to_process, total_chars


def _load_pdf(path: str, *, progress_callback: Optional[Callable[[dict], None]] = None) -> DocumentLoadResult:
    file_name = os.path.basename(path)
    load_start = time.perf_counter()
    native_text, meta = extract_pdf_text(path)
    meta = dict(meta)
    meta["loader"] = "pdf-tesseract"
    result = DocumentLoadResult(text="", metadata=meta)

    try:
        version = _ensure_tesseract_available()
        logger.info("Insertion PDF OCR tesseract version=%s file=%s", version, file_name)
        text, processed_pages, ocr_chars = _tesseract_ocr_pdf_file(
            path,
            max_pages=PDF_OCR_MAX_PAGES,
            zoom=PDF_OCR_ZOOM,
            progress_callback=progress_callback,
        )
        result.text = text
        result.metadata["ocr_pages_processed"] = processed_pages
        result.metadata["ocr_chars"] = ocr_chars
    except Exception as e:
        result.warnings.append(
            f"pytesseract PDF OCR failed: {e}. Install the native tesseract binary (macOS: brew install tesseract)."
        )
        logger.warning("Insertion PDF OCR failed file=%s error=%s", file_name, e)

    if not result.text.strip() and native_text.strip():
        result.text = native_text
        result.metadata["loader"] = "pdf-text-fallback"
        result.warnings.append("Used native PDF text because pytesseract returned no text.")
        logger.info(
            "Insertion PDF native text fallback file=%s chars=%d",
            file_name,
            len(native_text),
        )

    if not result.text.strip():
        # Last-resort compatibility with the existing vision LLM path.
        try:
            result.page_images = pdf_page_images(path)
            result.needs_vision = bool(result.page_images)
        except Exception as e:
            result.warnings.append(f"Could not render PDF pages for vision fallback: {e}")
            logger.warning("Insertion PDF vision fallback render failed file=%s error=%s", file_name, e)
    logger.info(
        "Insertion PDF load complete file=%s loader=%s chars=%d pages_total=%s ocr_pages=%s needs_vision=%s elapsed_s=%.2f",
        file_name,
        result.metadata.get("loader"),
        len(result.text),
        result.metadata.get("page_count"),
        result.metadata.get("ocr_pages_processed", 0),
        result.needs_vision,
        time.perf_counter() - load_start,
    )
    return result


def _load_docx(path: str) -> DocumentLoadResult:
    from docx import Document  # type: ignore

    doc = Document(path)
    parts: List[str] = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = "\t".join((cell.text or "").strip() for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)
    return DocumentLoadResult(text="\n".join(parts), metadata={"loader": "docx"})


def _load_doc_legacy(path: str) -> DocumentLoadResult:
    # .doc (binary) isn't well supported in pure Python. Try textract-free fallback: read as bytes
    # and look for ASCII runs. Users should convert to docx for best results.
    raw = open(path, "rb").read()
    text_chunks: List[str] = []
    current: List[int] = []
    for b in raw:
        if 32 <= b <= 126 or b in (9, 10, 13):
            current.append(b)
        else:
            if len(current) > 8:
                text_chunks.append(bytes(current).decode("ascii", errors="ignore"))
            current = []
    if len(current) > 8:
        text_chunks.append(bytes(current).decode("ascii", errors="ignore"))
    return DocumentLoadResult(
        text="\n".join(text_chunks),
        metadata={"loader": "doc-legacy"},
        warnings=["Legacy .doc: text extraction is best-effort. Convert to .docx for fidelity."],
    )


def _load_xlsx(path: str) -> DocumentLoadResult:
    from openpyxl import load_workbook  # type: ignore

    wb = load_workbook(filename=path, data_only=True, read_only=True)
    lines: List[str] = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"# Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(cell.strip() for cell in cells):
                lines.append("\t".join(cells))
        lines.append("")
    return DocumentLoadResult(text="\n".join(lines), metadata={"loader": "xlsx"})


def _load_pptx(path: str) -> DocumentLoadResult:
    from pptx import Presentation  # type: ignore

    prs = Presentation(path)
    slides: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        bits: List[str] = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        bits.append(text)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_text = "\t".join((cell.text or "").strip() for cell in row.cells)
                    if row_text.strip():
                        bits.append(row_text)
        slides.append("\n".join(bits))
    return DocumentLoadResult(text="\n\n".join(slides), metadata={"loader": "pptx"})


def _load_image(path: str) -> DocumentLoadResult:
    """Load image as PNG bytes and OCR it with pytesseract."""
    from PIL import Image  # type: ignore

    ext = os.path.splitext(path)[1].lower()
    warnings: List[str] = []
    if ext == ".heic":
        try:
            import pillow_heif  # type: ignore

            pillow_heif.register_heif_opener()
        except Exception:
            pass
    with Image.open(path) as img:
        img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        png_bytes = buf.getvalue()

    text = ""
    needs_vision = False
    try:
        text = _tesseract_ocr_png(png_bytes)
        logger.info(
            "Insertion image OCR complete file=%s chars=%d image_bytes=%d",
            os.path.basename(path),
            len(text),
            len(png_bytes),
        )
    except Exception as e:
        needs_vision = True
        warnings.append(
            f"pytesseract image OCR failed: {e}. Install the native tesseract binary (macOS: brew install tesseract)."
        )
        logger.warning("Insertion image OCR failed file=%s error=%s", os.path.basename(path), e)
    if not text.strip():
        needs_vision = True
        warnings.append("pytesseract returned no text for image; falling back to vision OCR.")
    return DocumentLoadResult(
        text=text,
        images=[png_bytes],
        metadata={"loader": "image-tesseract", "format": ext.lstrip(".")},
        needs_vision=needs_vision,
        warnings=warnings,
    )


# ---------- dispatcher ----------


_EXT_LOADERS = {
    ".txt": _load_txt_md,
    ".md": _load_txt_md,
    ".json": _load_json,
    ".html": _load_html,
    ".htm": _load_html,
    ".pdf": _load_pdf,
    ".docx": _load_docx,
    ".doc": _load_doc_legacy,
    ".xlsx": _load_xlsx,
    ".xls": _load_xlsx,
    ".pptx": _load_pptx,
    ".jpg": _load_image,
    ".jpeg": _load_image,
    ".png": _load_image,
    ".heic": _load_image,
}


def load_file_to_text(path: str, progress_callback: Optional[Callable[[dict], None]] = None) -> DocumentLoadResult:
    """Dispatch to the correct loader based on the file extension."""
    ext = os.path.splitext(path)[1].lower()
    if ext not in _EXT_LOADERS:
        raise ValueError(
            f"Unsupported file type: {ext or '(none)'}. Supported: {', '.join(sorted(_EXT_LOADERS))}"
        )
    try:
        if ext == ".pdf":
            return _load_pdf(path, progress_callback=progress_callback)
        return _EXT_LOADERS[ext](path)
    except ImportError as e:
        raise RuntimeError(
            f"Missing dependency for {ext}: {e}. Install requirements.txt in the backend venv."
        ) from e


def image_to_data_url(png_bytes: bytes) -> str:
    """Encode PNG bytes as a data URL for multimodal LLM inputs."""
    b64 = base64.b64encode(png_bytes).decode("ascii")
    return f"data:image/png;base64,{b64}"
